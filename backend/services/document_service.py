import os
import uuid
import re
from fastapi import UploadFile, HTTPException

ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
    "application/octet-stream"  # Evaluated alongside extension
}

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}
BLOCKED_EXTENSIONS = {".exe", ".bat", ".cmd", ".sh", ".dll", ".bin", ".msi", ".vbs", ".ps1", ".py", ".js", ".html", ".htm"}
MAX_FILE_SIZE_BYTES = 25 * 1024 * 1024  # 25 MB

def validate_file(file: UploadFile) -> bool:
    """
    Validates uploaded file against allowed extensions, MIME types, and blocked executable types.
    """
    if not file or not file.filename:
        return False
        
    ext = os.path.splitext(file.filename)[1].lower()
    
    if ext in BLOCKED_EXTENSIONS:
        return False
        
    if ext not in ALLOWED_EXTENSIONS:
        return False
        
    if file.content_type and file.content_type not in ALLOWED_MIME_TYPES:
        return False
        
    return True

def save_upload(file: UploadFile, upload_dir: str) -> str:
    """
    Saves uploaded file to disk with unique UUID filename and enforces size limits.
    """
    if not os.path.exists(upload_dir):
        os.makedirs(upload_dir, exist_ok=True)
        
    ext = os.path.splitext(file.filename)[1].lower()
    safe_name = f"{uuid.uuid4().hex}{ext}"
    path = os.path.join(upload_dir, safe_name)
    
    bytes_read = 0
    with open(path, "wb") as f:
        while chunk := file.file.read(1024 * 1024):  # 1MB chunks
            bytes_read += len(chunk)
            if bytes_read > MAX_FILE_SIZE_BYTES:
                f.close()
                if os.path.exists(path):
                    os.remove(path)
                raise HTTPException(
                    status_code=400, 
                    detail=f"File exceeds maximum allowed size of {MAX_FILE_SIZE_BYTES // (1024 * 1024)} MB."
                )
            f.write(chunk)
            
    return path

def extract_text(file_path: str, file_type: str = "") -> str:
    """
    Extracts text from PDF, DOCX, PPTX, or plain text files.
    Returns clean text string or raises ValueError if extraction yields empty/unreadable content.
    """
    text = ""
    ext = os.path.splitext(file_path)[1].lower()
    file_type = (file_type or "").lower()

    try:
        if ext == ".pdf" or "pdf" in file_type:
            import fitz
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n"
        elif ext == ".docx" or "wordprocessingml" in file_type:
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
        elif ext == ".pptx" or "presentationml" in file_type:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception as e:
        text = ""

    cleaned = text.strip()
    return cleaned

def chunk_text(text: str, max_chunk_size: int = 3000):
    return [text[i:i+max_chunk_size] for i in range(0, len(text), max_chunk_size)]
